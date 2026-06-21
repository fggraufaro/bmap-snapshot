"""
BMAP Snapshot Deck Builder
==========================
Generates a branded 6-slide PowerPoint deck for any bank in BMAP.
Uses python-pptx for native vector charts (crisp at any zoom).
Opens directly in PowerPoint — no Google Slides workaround needed.

Setup (one time):
    pip install python-pptx anthropic requests

Usage — single bank:
    python bmap_snapshot.py --inst_key bank_123

Usage — batch from CSV:
    python bmap_snapshot.py --csv banks.csv

CSV format (one bank per line):
    inst_key,optional_name
    bank_123,Country Bank for Savings
    bank_456

Output: BMAP_Snapshot_<BankName>_<date>.pptx  (or a folder for batch)
"""

import argparse
import csv
import io
import json
import hashlib
import os
import sys
import time
import urllib.request
from datetime import datetime
from pathlib import Path

# ── Third-party ────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import ChartData
    from pptx.oxml.ns import qn
    import pptx.oxml as pxml
    from lxml import etree
except ImportError:
    sys.exit("Run: pip install python-pptx")

try:
    import anthropic
except ImportError:
    sys.exit("Run: pip install anthropic")

try:
    import requests
except ImportError:
    sys.exit("Run: pip install requests")

# ═══════════════════════════════════════════════════════════════
# CONFIG — edit these
# ═══════════════════════════════════════════════════════════════
SUPA_URL  = "https://tuiiywphoynbmkxpoyps.supabase.co"
SUPA_KEY  = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InR1aWl5d3Bob3luYm1reHBveXBzIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTc0MDg0NTMsImV4cCI6MjA3Mjk4NDQ1M30.8-JAz4WQRE3Fi6uH7xiYNTns92g-nV1A9pbUvSK549M"
ANTH_KEY  = os.environ.get("ANTHROPIC_API_KEY", "")  # set env var or paste here
LOGO_URL  = "https://fggraufaro.github.io/bmap-tools/Verlocity-Logo.png"
OUT_DIR   = Path(".")

# ═══════════════════════════════════════════════════════════════
# BRAND
# ═══════════════════════════════════════════════════════════════
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY    = rgb("1A2332")
NAVY_SOFT = rgb("3D4D63")  # darker than GRAY3 — used for narrative body copy, keeps contrast on white
TEAL    = rgb("1D9E75")
AMBER   = rgb("F5A623")
WHITE   = rgb("FFFFFF")
GRAY1   = rgb("F5F5F2")
GRAY2   = rgb("E8E8E5")
GRAY3   = rgb("778899")
INVEST  = rgb("27500A"); INVEST_L = rgb("EAF3DE")
ANALYZE = rgb("185FA5"); ANALYZE_L = rgb("E6F1FB")
DEFEND  = rgb("854F0B"); DEFEND_L  = rgb("FFF3E0")
JUSTIFY = rgb("A32D2D"); JUSTIFY_L = rgb("FCEBEB")

ZONE_C = {"Invest": INVEST, "Analyze": ANALYZE, "Defend": DEFEND, "Justify": JUSTIFY}
ZONE_L = {"Invest": INVEST_L, "Analyze": ANALYZE_L, "Defend": DEFEND_L, "Justify": JUSTIFY_L}

# Slide dimensions: 10" × 5.625" (widescreen)
W = Inches(10)
H = Inches(5.625)

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════
SCHEMA_MAP = {
    'branch_opportunity_base':      'analytics',
    'bank_financial_snapshot_latest': 'analytics',
    'branch_target_competitors':    'analytics',
    'dim_institutions':             'ref',
    'bank_website':                 'ref',
}

def supabase(table, params):
    url = f"{SUPA_URL}/rest/v1/{table}?{params}"
    schema = SCHEMA_MAP.get(table, 'public')
    headers = {"apikey": SUPA_KEY, "Authorization": f"Bearer {SUPA_KEY}"}
    if schema != 'public':
        headers['Accept-Profile'] = schema
    r = requests.get(url, headers=headers, timeout=30)
    r.raise_for_status()
    return r.json()



# ═══════════════════════════════════════════════════════════════
# PERSONA GENERATION
# ═══════════════════════════════════════════════════════════════

def _branch_fingerprint(br):
    """Hash of sorted branch IDs + zones — fingerprint for this branch selection."""
    ids = sorted(str(b.get("uninumbr","")) for b in br)
    zones = sorted(str(b.get("opportunity_zone","")) for b in br)
    raw = "|".join(ids) + "||" + "|".join(zones)
    return hashlib.md5(raw.encode()).hexdigest()[:16]

def _branch_zone_summary(br):
    """e.g. Invest:2,Analyze:4,Defend:1,Justify:3"""
    from collections import Counter
    c = Counter(b.get("opportunity_zone","Unknown") for b in br)
    return ",".join(f"{z}:{c[z]}" for z in ["Invest","Analyze","Defend","Justify"] if z in c)

def fetch_or_generate_personas(ik, institution_name, br, data):
    """
    Check Supabase for approved/draft personas for this inst_key + branch_fingerprint.
    If found -> return existing.
    If not found -> generate with Claude + web search -> save as draft -> return.
    """
    fingerprint = _branch_fingerprint(br)
    zone_summary = _branch_zone_summary(br)

    # 1. Check for existing approved or draft personas
    url = (f"{SUPA_URL}/rest/v1/persona_runs"
           f"?inst_key=eq.{ik}&branch_fingerprint=eq.{fingerprint}"
           f"&status=in.(approved,draft)&order=run_date.desc&limit=1")
    r = requests.get(url, headers={"apikey": SUPA_KEY, "Authorization": f"Bearer {SUPA_KEY}"}, timeout=15)
    rows = r.json() if r.ok else []

    if rows:
        p = rows[0]
        print(f"  Using existing personas (status: {p.get('status')}, run: {p.get('run_date','')[:10]})")
        return _parse_persona_run(p)

    # 2. No existing — generate new
    print("  Generating new personas with Claude + web search...")
    personas = _generate_personas(ik, institution_name, br, data)
    if not personas:
        # Fallback: build 3 generic personas from branch demographics so the
        # slide is never silently skipped when Claude generation fails.
        rows_data = data.get("rows", [])
        metro = rows_data[0].get("metro", institution_name) if rows_data else institution_name
        avg_income = sum(float(r.get("household_income") or 0) for r in rows_data) / max(len(rows_data), 1)
        personas = [
            {
                "name": f"The {metro} Rate Seeker",
                "age": "35–54",
                "income": f"${avg_income/1000:.0f}k avg household" if avg_income else "Mid-income",
                "occupation": "Professional / Dual-income household",
                "insight": f"Actively comparing CD and savings rates in {metro} — rate environment is top of mind.",
                "moment": "CD ladder to lock in rates before the next Fed move",
                "why_now": "Fed rate uncertainty is driving depositors to act now rather than wait.",
            },
            {
                "name": f"The {metro} Community Loyalist",
                "age": "45–65",
                "income": f"${avg_income/1000:.0f}k avg household" if avg_income else "Mid-income",
                "occupation": "Small business owner / Pre-retiree",
                "insight": "Values local relationships and trust — open to consolidating accounts at one institution.",
                "moment": "Primary checking anchor + savings relationship deepening",
                "why_now": "National bank frustration is at a high — community banking is a differentiated story.",
            },
            {
                "name": f"The {metro} Digital Switcher",
                "age": "28–44",
                "income": f"${max(50, avg_income/1000 - 15):.0f}k avg household" if avg_income else "Mid-income",
                "occupation": "Tech-adjacent professional / Young family",
                "insight": "Mobile-first, rate-aware — looking for a better savings rate without sacrificing digital experience.",
                "moment": "High-yield savings account as entry product → CD conversion",
                "why_now": "Online bank rates are plateauing — community banks with competitive rates can win on trust.",
            },
        ]
        print(f"  Using fallback personas for {institution_name}")

    # 3. Save as draft to Supabase
    sf = lambda v: float(v) if v is not None else None
    rows_data = data.get("rows", [])
    avg_income = sum(float(r.get("household_income") or 0) for r in rows_data) / max(len(rows_data),1)
    avg_pop_growth = sum(float(r.get("yoy_pop_growth") or 0) for r in rows_data) / max(len(rows_data),1) * 100
    avg_zhvi = sum(float(r.get("zhvi_yoy_pct") or 0) for r in rows_data) / max(len(rows_data),1)
    metro = rows_data[0].get("metro","") if rows_data else ""

    payload = {
        "inst_key":           ik,
        "institution_name":   institution_name,
        "branch_fingerprint": fingerprint,
        "branch_count":       len(br),
        "branch_zones":       zone_summary,
        "metro":              metro,
        "avg_income":         round(avg_income, 0) if avg_income else None,
        "pop_growth_pct":     round(avg_pop_growth, 2) if avg_pop_growth else None,
        "zhvi_yoy_pct":       round(avg_zhvi, 2) if avg_zhvi else None,
        "status":             "draft",
        "p1_name":    personas[0].get("name"),
        "p1_age":     personas[0].get("age"),
        "p1_income":  personas[0].get("income"),
        "p1_occupation": personas[0].get("occupation"),
        "p1_insight": personas[0].get("insight"),
        "p1_moment":  personas[0].get("moment"),
        "p1_why_now": personas[0].get("why_now"),
        "p2_name":    personas[1].get("name") if len(personas)>1 else None,
        "p2_age":     personas[1].get("age") if len(personas)>1 else None,
        "p2_income":  personas[1].get("income") if len(personas)>1 else None,
        "p2_occupation": personas[1].get("occupation") if len(personas)>1 else None,
        "p2_insight": personas[1].get("insight") if len(personas)>1 else None,
        "p2_moment":  personas[1].get("moment") if len(personas)>1 else None,
        "p2_why_now": personas[1].get("why_now") if len(personas)>1 else None,
        "p3_name":    personas[2].get("name") if len(personas)>2 else None,
        "p3_age":     personas[2].get("age") if len(personas)>2 else None,
        "p3_income":  personas[2].get("income") if len(personas)>2 else None,
        "p3_occupation": personas[2].get("occupation") if len(personas)>2 else None,
        "p3_insight": personas[2].get("insight") if len(personas)>2 else None,
        "p3_moment":  personas[2].get("moment") if len(personas)>2 else None,
        "p3_why_now": personas[2].get("why_now") if len(personas)>2 else None,
    }

    save_url = f"{SUPA_URL}/rest/v1/persona_runs"
    save_r = requests.post(save_url,
        headers={"apikey": SUPA_KEY, "Authorization": f"Bearer {SUPA_KEY}",
                 "Content-Type": "application/json", "Prefer": "return=minimal"},
        json=payload, timeout=15)
    if save_r.ok:
        print(f"  Personas saved as draft (fingerprint: {fingerprint})")
    else:
        print(f"  Warning: could not save personas ({save_r.status_code})")

    return personas


def _parse_persona_run(p):
    """Convert a persona_runs DB row into the personas list format."""
    personas = []
    for prefix in ["p1","p2","p3"]:
        if p.get(f"{prefix}_name"):
            personas.append({
                "name":       p.get(f"{prefix}_name",""),
                "age":        p.get(f"{prefix}_age",""),
                "income":     p.get(f"{prefix}_income",""),
                "occupation": p.get(f"{prefix}_occupation",""),
                "insight":    p.get(f"{prefix}_insight",""),
                "moment":     p.get(f"{prefix}_moment",""),
                "why_now":    p.get(f"{prefix}_why_now",""),
            })
    return personas if personas else None


def _generate_personas(ik, institution_name, br, data):
    """Generate 3 personas using Claude with web search, grounded in branch demographics."""
    if not ANTH_KEY:
        return None

    rows = data.get("rows", [])
    # Only use Invest + Analyze branches for persona targeting
    target_br = [b for b in rows if b.get("opportunity_zone") in ("Invest","Analyze")
                 and float(b.get("latest_dep") or 0) >= 5e6]
    if not target_br:
        target_br = [b for b in rows if float(b.get("latest_dep") or 0) >= 5e6]

    # Build demographic context from branch ZIPs
    def safe_avg(key):
        vals = [float(v) for b in target_br if (v:=b.get(key)) not in (None,"")]
        return round(sum(vals)/len(vals), 2) if vals else None

    avg_income    = safe_avg("household_income")
    avg_inc_yoy   = safe_avg("yoy_income_growth")
    avg_pop_yoy   = safe_avg("yoy_pop_growth")
    avg_zhvi_yoy  = safe_avg("zhvi_yoy_pct")
    avg_dep_yoy   = safe_avg("yoy_deposits")
    metro         = target_br[0].get("metro","") if target_br else ""
    states        = list(set(b.get("stalpbr","") for b in target_br))
    cities        = list(set(b.get("citybr","") for b in target_br))[:5]

    demo_ctx = f"""Institution: {institution_name}
Market: {metro} | States: {", ".join(states)} | Key cities: {", ".join(cities)}
Branch pool: {len(target_br)} Invest/Analyze branches (>$5M deposits)

DEMOGRAPHIC SIGNALS (Census ACS 2024, averaged across target branch ZIPs):
- Avg household income: {f"${avg_income:,.0f}" if avg_income else "N/A"}
- Income growth YoY: {f"+{avg_inc_yoy*100:.1f}%" if avg_inc_yoy else "N/A"}
- Population growth YoY: {f"+{avg_pop_yoy*100:.1f}%" if avg_pop_yoy else "N/A"}
- Home value appreciation (ZHVI YoY): {f"+{avg_zhvi_yoy:.1f}%" if avg_zhvi_yoy else "N/A"}
- Branch deposit growth vs peers: {f"+{avg_dep_yoy*100:.1f}%" if avg_dep_yoy else "N/A"}"""

    system = """You are Verlocity's Audience Intelligence Director. Your job is to identify the 3 most valuable deposit-growth personas for a community bank based on their specific market demographics and current economic context.

PERSONA PHILOSOPHY:
- These are real people, not demographic buckets. Give them names that feel human ("The Worcester Accumulator", "The Gulf Coast Professional")  
- Ground every insight in the actual data provided
- Use web search to add current market context — what's happening economically in this specific metro that creates a banking moment RIGHT NOW
- Banking moments must be specific and actionable — not "savings account" but "CD ladder to lock in rates before Fed cuts"
- Why Now must reference something current and real — a rate environment shift, local employer, demographic wave

Return ONLY a JSON array of exactly 3 objects, each with:
{
  "name": "The [Descriptor] [Type]",
  "age": "35-54",
  "income": "$87K avg household",
  "occupation": "Primary occupation cluster",
  "insight": "One sentence — who they are and what defines their financial life right now",
  "moment": "Specific banking product/behavior opportunity",
  "why_now": "One sentence — why this is the moment to reach them, referencing current market conditions"
}

No markdown, no explanation, ONLY the JSON array."""

    import anthropic
    client = anthropic.Anthropic(api_key=ANTH_KEY)

    try:
        messages = [{
            "role": "user",
            "content": f"""Generate 3 target personas for {institution_name} based on this data:

{demo_ctx}

Search for current economic conditions, major employers, demographic trends, and banking behavior in {metro} to enrich the personas with real market context. Then return ONLY the JSON array — no explanation, no markdown."""
        }]

        # Agentic loop — handle web_search tool use turns properly
        txt = ""
        for attempt in range(6):  # max 6 turns (search calls)
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2000,
                tools=[{"type": "web_search_20250305", "name": "web_search"}],
                system=system,
                messages=messages
            )

            # Collect any text from this turn
            for block in response.content:
                if hasattr(block, "text") and block.text:
                    txt += block.text

            # If model is done (end_turn or no tool calls), break
            if response.stop_reason == "end_turn":
                break

            # If tool_use, feed results back and continue
            tool_calls = [b for b in response.content if b.type == "tool_use"]
            if not tool_calls:
                break

            # Append assistant turn + tool results to messages
            messages.append({"role": "assistant", "content": response.content})
            tool_results = []
            for tc in tool_calls:
                # web_search results are returned by the API automatically in the next turn
                tool_results.append({
                    "type": "tool_result",
                    "tool_use_id": tc.id,
                    "content": "Search completed."
                })
            messages.append({"role": "user", "content": tool_results})

        # Parse JSON from collected text
        clean = txt.replace("```json", "").replace("```", "").strip()
        s = clean.find("[")
        e = clean.rfind("]") + 1
        if s >= 0 and e > s:
            personas = json.loads(clean[s:e])
            if personas and len(personas) >= 1:
                print(f"  ✓ Generated {len(personas)} personas for {institution_name}")
                return personas[:3]

        print(f"  Could not parse personas JSON for {institution_name}. Response: {txt[:200]}")
        return None

    except Exception as e:
        print(f"  Persona generation error for {institution_name}: {e}")
        return None


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=11, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, italic=False, font="Calibri", valign="top",
             shrink_to_fit=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if shrink_to_fit:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM,
                           "center": MSO_ANCHOR.MIDDLE}.get(valign, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def truncate_label(text, max_len):
    """Truncate to max_len, breaking on the last word boundary rather than mid-word."""
    if len(text) <= max_len:
        return text
    cut = text[:max_len]
    if " " in cut:
        cut = cut[:cut.rindex(" ")]
    return cut.rstrip() + "…"

def add_chrome(slide, page_num, label, logo_bytes):
    """Verlocity sidebar + logo + section label + page number"""
    add_rect(slide, 0, 0, 0.28, 5.625, NAVY)
    add_rect(slide, 0.28, 0, 0.08, 5.625, TEAL)
    if logo_bytes:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(0.42), Inches(5.04), Inches(1.55), Inches(0.36))
    if label:
        pill = add_rect(slide, 8.16, 0.14, 1.76, 0.30, TEAL)
        add_text(slide, label, 8.16, 0.14, 1.76, 0.30, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, str(page_num), 9.50, 5.28, 0.38, 0.20, size=9, color=GRAY3, align=PP_ALIGN.RIGHT)

def add_narrative(slide, n, y0):
    """Left narrative column: headline / rule / spoken / bullets / close bar"""
    add_text(slide, n.get("headline",""), 0.45, y0, 5.6, 0.78,
             size=24, bold=True, color=NAVY, valign="bottom")
    add_rect(slide, 0.45, y0+0.84, 5.6, 0.04, TEAL)
    # Spoken line: taller box + shrink-to-fit so longer AI text never
    # overlaps the bullets below it instead of just clipping invisibly.
    add_text(slide, n.get("spoken",""), 0.45, y0+0.96, 5.6, 0.62,
             size=9.5, italic=True, color=NAVY_SOFT, shrink_to_fit=True)
    bullets = n.get("bullets", [])
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(y0+1.66), Inches(5.6), Inches(1.92))
        tf = tb.text_frame; tf.word_wrap = True
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(9.5); p.font.color.rgb = NAVY; p.font.name = "Calibri"
            p.space_after = Pt(5)
    add_rect(slide, 0.45, y0+3.68, 5.6, 0.46, NAVY)
    add_text(slide, n.get("close",""), 0.56, y0+3.68, 5.4, 0.46,
             size=9.5, bold=True, color=WHITE, valign="center", shrink_to_fit=True)

def fetch_logo():
    try:
        with urllib.request.urlopen(LOGO_URL, timeout=10) as r:
            return r.read()
    except Exception as e:
        print(f"  ⚠  Logo fetch failed ({e}) — slides will have no logo")
        return None

# ═══════════════════════════════════════════════════════════════
# DATA FETCH
# ═══════════════════════════════════════════════════════════════
def _fetch_brokered(ik):
    """
    Fetch brokered deposit data from raw_schedule_RCE + raw_schedule_RC.
    Returns dict with brokered metrics, or None if below 15% threshold.
    Uses RCON2365 (brokered deposits) / RCON2385 (total deposits RCE).
    All values in $thousands from FFIEC call report.
    """
    THRESHOLD = 0.15   # 15% — meaningful brokered pressure

    # Extract RSSDID from inst_key (bank_463735 → 463735)
    rssdid = ik.replace("bank_", "").replace("cu_", "")

    # Use Supabase RPC-style SQL via the raw PostgREST endpoint
    # We'll query raw_schedule_RCE directly
    url = (
        f"{SUPA_URL}/rest/v1/raw_schedule_RCE"
        f'?IDRSSD=eq.{rssdid}'
        f'&period=eq.2025-12-31'
        f'&select=RCON2365,RCON2385,RCON0352,RCON5590'
        f'&limit=1'
    )
    import requests as _req
    try:
        r = _req.get(url,
            headers={"apikey": SUPA_KEY, "Authorization": f"Bearer {SUPA_KEY}", "Accept-Profile": "raw"},
            timeout=15)
        rows = r.json()
    except Exception as e:
        print(f"  [brokered] fetch error: {e}")
        return None

    if not rows:
        return None

    row = rows[0]
    try:
        brokered  = float(row.get("RCON2365") or 0)   # $thousands
        total_dep = float(row.get("RCON2385") or 0)   # $thousands
        savings   = float(row.get("RCON0352") or 0)   # $thousands
        cds       = float(row.get("RCON5590") or 0)   # $thousands
    except (TypeError, ValueError):
        return None

    if total_dep == 0:
        return None

    brok_pct = brokered / total_dep

    if brok_pct < THRESHOLD:
        return None   # below threshold — don't show module

    return {
        "brokered_M":   round(brokered   / 1000, 1),   # convert to $M
        "total_dep_M":  round(total_dep  / 1000, 1),
        "savings_M":    round(savings    / 1000, 1),
        "cds_M":        round(cds        / 1000, 1),
        "brokered_pct": round(brok_pct * 100, 1),
    }


def fetch_bank_data(ik):
    print(f"  Fetching branch data...")
    rows = supabase("branch_opportunity_base",
        f"inst_key=eq.{ik}&select=namefull,latest_dep,yoy_deposits,avg_comp_yoy,"
        "opportunity_score,opportunity_zone,market_growth_score,inv_density_norm_winsor")

    print(f"  Fetching branch details...")
    br = supabase("branch_opportunity_base",
        f"inst_key=eq.{ik}&select=uninumbr,namebr,citybr,stalpbr,latest_dep,"
        "yoy_deposits,opportunity_score,opportunity_zone,matrix_quadrant,priority_tier,campaign"
        "&order=opportunity_score.desc&limit=50")

    print(f"  Fetching network target...")
    tgt_arr = supabase("vw_network_top_targets",
        f"my_inst_key=eq.{ik}&select=target_institution,branches_in_radius,"
        "avg_vuln_score,avg_yoy_pct,target_roa,target_efficiency_ratio,dominant_zone"
        "&order=network_rank.asc&limit=1")

    print(f"  Fetching financials...")
    fin_arr = supabase("bank_financial_snapshot_latest",
        f"inst_key=eq.{ik}&select=*&limit=1")

    print(f"  Fetching brokered deposits...")
    brok = _fetch_brokered(ik)

    return {
        "ik":       ik,
        "bankName": rows[0].get("namefull", ik) if rows else ik,
        "rows":     rows or [],
        "br":       br   or [],
        "tgt":      tgt_arr[0] if tgt_arr else None,
        "fin":      fin_arr[0] if fin_arr else {},
        "brokered": brok,
    }

# ═══════════════════════════════════════════════════════════════
# AI NARRATIVES
# ═══════════════════════════════════════════════════════════════
def get_narratives(data):
    if not ANTH_KEY:
        print("  ⚠  No ANTHROPIC_API_KEY — using placeholder narratives")
        empty = {"headline":"","spoken":"","bullets":[],"close":""}
        return {k: empty for k in ["network","priority","financial","nextsteps"]}

    rows = data["rows"]; fin = data["fin"]; tgt = data["tgt"]
    br   = data["br"];   bankName = data["bankName"]

    def sf(v, default=0):
        try:
            return float(v) if v is not None else default
        except (TypeError, ValueError):
            return default

    tot      = sum(sf(r.get("latest_dep")) for r in rows)
    avg      = lambda v: sum(sf(r.get(v)) for r in rows) / max(len(rows),1)
    invest   = sum(1 for r in rows if r.get("opportunity_zone")=="Invest")
    analyze  = sum(1 for r in rows if r.get("opportunity_zone")=="Analyze")
    defend   = sum(1 for r in rows if r.get("opportunity_zone")=="Defend")
    justify  = sum(1 for r in rows if r.get("opportunity_zone")=="Justify")
    bankYoY  = avg("yoy_deposits")*100
    compYoY  = avg("avg_comp_yoy")*100
    gap      = bankYoY - compYoY
    avgScore = avg("opportunity_score")

    sig_br = [b for b in br if sf(b.get("latest_dep")) >= 5e6]
    top3 = sorted(sig_br or br, key=lambda b: sf(b.get("opportunity_score")), reverse=True)[:3]
    top3_str = "; ".join(
        f"{b['namebr'].split('--')[-1].strip()} "
        f"(${sf(b.get('latest_dep'))/1e6:.0f}M, "
        f"{sf(b.get('yoy_deposits'))*100:.1f}% YoY, "
        f"score {sf(b.get('opportunity_score')):.0f})"
        for b in top3
    )

    comp_str = (
        f"Primary competitor — {tgt['branches_in_radius']} overlap branches, "
        f"avg vulnerability score {sf(tgt.get('avg_vuln_score')):.0f}/100, "
        f"their YoY {sf(tgt.get('avg_yoy_pct')):.1f}%"
    ) if tgt else "N/A"

    brok = data.get("brokered")
    brok_str = (
        f"Brokered deposits: {brok['brokered_pct']}% of total "
        f"(${brok['brokered_M']:.0f}M of ${brok['total_dep_M']:.0f}M total). "
        f"Direct savings base: ${brok['savings_M']:.0f}M. "
        f"Conversion opportunity: replace expensive brokered funding with "
        f"direct customer deposits via savings→CD funnel strategy."
    ) if brok else "Brokered deposits: not a material factor for this bank."

    ctx = f"""Bank: {bankName} | {len(rows)} branches | ${tot/1e9:.2f}B deposits
Deposit YoY: +{bankYoY:.1f}% | Peer avg (competitor growth in bank\'s own markets): +{compYoY:.1f}% | Gap: {gap:+.1f}pp
Avg opp score: {avgScore:.1f}/100 | Zones: Invest {invest} | Analyze {analyze} | Defend {defend} | Justify {justify}
ROA: {fin.get('roa','—')}% | NIM: {fin.get('nim','—')}% | Efficiency: {fin.get('efficiency_ratio','—')}%
Net income YoY: {fin.get('net_income_yoy_pct','—')}% | Tier 1: {fin.get('tier1_capital_pct','—')}%
Primary competitor: {comp_str}
Brokered deposit situation: {brok_str}
Top 3 branches: {top3_str}"""

    system = """You are BMAP Executive Strategist at Verlocity Princeton Partners Group. Write boardroom-quality slide narratives grounded in exact numbers. These slides open conversations — they inform, they don't pressure. Never prescribe specific actions the bank can execute without Verlocity. Name the opportunity clearly and let the data make the case; avoid hype language, urgency gimmicks, or anything that reads as a sales pitch rather than an analysis.

VERLOCITY PLATFORM CONTEXT (reference naturally where relevant):
Verlocity's marketing intelligence platform is built around four capabilities, all live today:
- BMAP — Market Truth: branch-level scoring across deposits, growth, and competitive density; where to invest, defend, and exit.
- AudienceFinder & MediaPredict — Media Performance: precision-targeted depositor campaigns built on the audiences BMAP identifies, converting savings openers into sticky CD relationships.
- Omnibranch — Branch Orchestration: unifies deposit data, market signals, and campaign execution into one operating view across the network.
- Predictive ROI — Accountability: forecasts the return on marketing spend before it happens, then tracks performance against that forecast.
The BMAP Snapshot is the entry point — it's the data layer the other three capabilities act on.
When writing nextsteps, position Verlocity as the partner that turns BMAP's market read into targeted campaigns (AudienceFinder/MediaPredict), coordinated execution (Omnibranch), and measurable return (Predictive ROI). Reference the platform naturally — the bank is seeing BMAP; help them understand what it connects to.
WinShare model: Verlocity invests alongside the bank — compensation tied to results, not effort.

BROKERED DEPOSIT CONTEXT (use if brokered pressure is present):
Brokered deposits are expensive, rate-sensitive, and leave when a better rate appears. The savings→CD funnel strategy converts brokered volume to sticky direct customer deposits at lower cost. If brokered pressure is present, frame it as a funded account conversion opportunity for AudienceFinder and MediaPredict campaigns targeting savings account openers who convert to CDs.

IMPORTANT RULES:
- Never name competitors directly. Refer to them as "your primary competitor" or "a regional competitor".
- Never prescribe specific budget reallocation amounts or staff actions.
- Bullets reveal WHAT the data shows, not HOW to fix it.
- Close bars end with a genuine question or invitation to discuss — never a pressure tactic or implied cost of inaction.
- "Peer avg" means competitor deposit growth in the bank's own markets.
- nextsteps headline: make it about what the Verlocity platform delivers as a whole, not just BMAP.
- nextsteps spoken: 300-330 characters max, factual tone. Note that BMAP is the data foundation the other three capabilities build on.
- nextsteps bullets: each bullet names one of the four capabilities and the concrete outcome it drives.
- nextsteps close: invite a conversation about next steps — direct and respectful, not a hard sell.

Return ONLY valid JSON — no markdown, no explanation:
{"slides":[
  {"id":"network","headline":"strong claim max 9 words","spoken":"ONE sentence, 300-330 characters max. The single sharpest number and comparison Tom would say walking in — do not stack multiple metrics. Factual tone, must read naturally aloud in one breath.","bullets":["data point with number","competitive threat — no competitor name","gap or risk that needs Verlocity to solve"],"close":"invitation to dig deeper — not a prescription"},
  {"id":"priority","headline":"...","spoken":"...","bullets":[...],"close":"..."},
  {"id":"financial","headline":"...","spoken":"...","bullets":[...],"close":"..."},
  {"id":"nextsteps","headline":"...","spoken":"...","bullets":[...],"close":"..."}
]}"""

    print(f"  Generating AI narratives...")
    client = anthropic.Anthropic(api_key=ANTH_KEY)
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2400,
        system=system,
        messages=[{"role":"user","content":f"Generate 4-slide narratives for:\n\n{ctx}"}]
    )
    raw = msg.content[0].text.strip().replace("```json","").replace("```","").strip()
    narr = {}
    try:
        for s in json.loads(raw)["slides"]:
            narr[s["id"]] = s
    except Exception as e:
        print(f"  ⚠  JSON parse failed ({e}) — using empty narratives")
    N = lambda k: narr.get(k, {"headline":"","spoken":"","bullets":[],"close":""})
    result = {k: N(k) for k in ["network","priority","financial","nextsteps"]}

    # Safety net: the "spoken" box is 5.6"w x 0.62"h at 9.5pt italic Calibri,
    # which holds ~330 characters before autofit shrinking hits its practical
    # floor and text starts overflowing instead of shrinking. The prompt asks
    # for 300-330 chars, but the model doesn't always follow that exactly —
    # truncate here so a box never receives more than it can actually hold,
    # regardless of what comes back from the API.
    SPOKEN_CHAR_LIMIT = 330
    for k in result:
        spoken = result[k].get("spoken", "")
        if len(spoken) > SPOKEN_CHAR_LIMIT:
            print(f"  ⚠  '{k}' spoken line was {len(spoken)} chars — truncating to {SPOKEN_CHAR_LIMIT}")
            result[k]["spoken"] = truncate_label(spoken, SPOKEN_CHAR_LIMIT)

    return result

# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_cover(prs, d, logo_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title and slide.shapes.title.element.getparent().remove(slide.shapes.title.element)

    add_rect(slide, 0, 0, 0.28, 5.625, NAVY)
    add_rect(slide, 0.28, 0, 0.08, 5.625, TEAL)

    # Logo top
    if logo_bytes:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(1.2), Inches(0.28), Inches(2.45), Inches(0.54))

    # Teal rule
    add_rect(slide, 1.2, 0.96, 8.6, 0.05, TEAL)

    # Bank name
    add_text(slide, d["bankName"], 1.2, 1.1, 8.5, 0.88,
             size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_text(slide, "BMAP Market Snapshot", 1.2, 2.04, 8.5, 0.34,
             size=16, color=NAVY)
    add_text(slide, d["date"], 1.2, 2.42, 8.5, 0.26, size=11, color=GRAY3)

    # 4 KPI tiles
    kpis = [
        (d["branchCount"],  "BRANCHES"),
        (d["deposits"],     "TOTAL DEPOSITS"),
        (d["avgScore"],     "AVG OPP SCORE"),
        (d["gap"],          "VS PEER AVG"),
    ]
    gap_neg = d["gapNeg"]
    for i, (val, lbl) in enumerate(kpis):
        kx = 1.2 + i*2.14
        add_rect(slide, kx, 2.82, 2.0, 0.88, GRAY1, GRAY2, Pt(0.5))
        c = JUSTIFY if (i==3 and gap_neg) else NAVY
        add_text(slide, val, kx, 2.88, 2.0, 0.48,
                 size=22, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, kx, 3.34, 2.0, 0.26,
                 size=7, bold=True, color=GRAY3, align=PP_ALIGN.CENTER)

    # 4 zone tiles
    zones = [
        (str(d["invest"]),  "INVEST",   INVEST,  INVEST_L),
        (str(d["analyze"]), "ANALYZE",  ANALYZE, ANALYZE_L),
        (str(d["defend"]),  "DEFEND",   DEFEND,  DEFEND_L),
        (str(d["justify"]), "JUSTIFY",  JUSTIFY, JUSTIFY_L),
    ]
    for i, (val, lbl, c, bg) in enumerate(zones):
        zx = 1.2 + i*2.14
        add_rect(slide, zx, 3.82, 2.0, 0.72, bg, c, Pt(0.8))
        add_text(slide, val, zx, 3.86, 2.0, 0.36, size=20, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, zx, 4.22, 2.0, 0.24, size=8,  bold=True, color=c, align=PP_ALIGN.CENTER)

    # Logo bottom-left + footer
    if logo_bytes:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(0.42), Inches(5.04), Inches(1.55), Inches(0.36))
    add_text(slide,
        f"Confidential  ·  Verlocity Princeton Partners Group  ·  {datetime.now().year}",
        1.2, 5.3, 8.5, 0.2, size=7.5, color=GRAY3, align=PP_ALIGN.CENTER)


def build_network(prs, d, narr, logo_bytes, page_num=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_num, "MARKET OVERVIEW", logo_bytes)
    add_narrative(slide, narr["network"], 0.14)

    # 4 KPI tiles right
    kpis = [
        (d["deposits"],   "TOTAL DEPOSITS", GRAY1,    NAVY),
        (d["avgScore"],   "AVG OPP SCORE",  GRAY1,    NAVY),
        (d["depositYoY"], "DEPOSIT YoY",    INVEST_L if not d["gapNeg"] else JUSTIFY_L,
                                            INVEST   if not d["gapNeg"] else JUSTIFY),
        (d["gap"],        "GAP VS MKT PEERS",   INVEST_L if not d["gapNeg"] else JUSTIFY_L,
                                            INVEST   if not d["gapNeg"] else JUSTIFY),
    ]
    for i, (val, lbl, bg, vc) in enumerate(kpis):
        kx = 6.46 + (i%2)*1.76
        ky = 0.14 + (i//2)*0.90
        add_rect(slide, kx, ky, 1.62, 0.78, bg, GRAY2, Pt(0.4))
        add_text(slide, val, kx, ky+0.06, 1.62, 0.42, size=19, bold=True, color=vc, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, kx, ky+0.52, 1.62, 0.20, size=6.5, bold=True, color=GRAY3, align=PP_ALIGN.CENTER)

    # ── ZONE CHART — horizontal paired bars, % branches vs % deposits ──
    # Why this replaces the table: a table that lists every zone works at
    # 4 branches but says nothing useful at 100+ branches, and it duplicates
    # facts the narrative paragraph already states. This chart instead
    # surfaces the thing a table can't show at a glance: the MISMATCH between
    # how many branches sit in a zone and how much money sits there — e.g.
    # a zone with few branches but most of the deposits, or vice versa. That
    # mismatch is the actual strategic signal, and it reads the same way
    # whether the bank has 4 branches or 400, because both series are
    # normalized to % of total rather than raw counts/dollars.
    #
    # Native XL_CHART_TYPE (not a hand-built donut/table) — the donut chart
    # tried earlier failed in real PowerPoint (see note on the old table
    # below); BAR_CLUSTERED is a standard chart type with no custom XML
    # patches needed, so it renders identically everywhere.
    add_text(slide, "Zones", 6.0, 1.92, 3.8, 0.30, size=14, bold=True, color=NAVY)

    ZONE_VIVID = {
        "Invest":  (INVEST,  d["invest"],  d["depInvest"]),
        "Analyze": (ANALYZE, d["analyze"], d["depAnalyze"]),
        "Defend":  (DEFEND,  d["defend"],  d["depDefend"]),
        "Justify": (JUSTIFY, d["justify"], d["depJustify"]),
    }
    total_branches  = sum(v[1] for v in ZONE_VIVID.values()) or 1
    total_dep_zones = sum(v[2] for v in ZONE_VIVID.values()) or 1

    # Order top-to-bottom as Invest/Analyze/Defend/Justify; chart_data
    # categories render bottom-to-top in a horizontal bar by default, so
    # reverse the order going in to get Invest at the top visually.
    zone_order = ["Justify", "Defend", "Analyze", "Invest"]
    pct_branches = [round(ZONE_VIVID[z][1] / total_branches * 100, 1) for z in zone_order]
    pct_deposits = [round(ZONE_VIVID[z][2] / total_dep_zones * 100, 1) for z in zone_order]

    chart_data = ChartData()
    chart_data.categories = zone_order
    chart_data.add_series("% of Branches", pct_branches)
    chart_data.add_series("% of Deposits", pct_deposits)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(6.0), Inches(2.26), Inches(3.8), Inches(2.55),
        chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = GRAY3
    chart.has_title = False

    # Branch-count series: solid zone color. Deposit series: lighter tint
    # of the same zone color per category — done per-point since each
    # category needs its own zone color, not one color per series.
    series_branches, series_deposits = chart.series[0], chart.series[1]
    for i, z in enumerate(zone_order):
        c = ZONE_C[z]
        series_branches.points[i].format.fill.solid()
        series_branches.points[i].format.fill.fore_color.rgb = c
        series_branches.points[i].format.line.fill.background()

        cl = ZONE_L[z]
        series_deposits.points[i].format.fill.solid()
        series_deposits.points[i].format.fill.fore_color.rgb = cl
        series_deposits.points[i].format.line.fill.background()

    va = chart.value_axis
    va.tick_labels.font.size = Pt(7.5)
    va.tick_labels.font.color.rgb = GRAY3
    va.has_major_gridlines = False
    va.minimum_scale = 0

    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(9.5)
    ca.tick_labels.font.color.rgb = NAVY
    ca.tick_labels.font.bold = True

    # Plot-level data labels (works across series; per-point dLbl XML like
    # the old donut attempt is what broke in real PowerPoint, so this sticks
    # to chart-level label settings only)
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0"%"'
    dl.number_format_is_linked = False
    dl.font.size = Pt(7.5)
    dl.font.color.rgb = NAVY

    # Small caption under the chart with the raw totals, so the % figures
    # have an anchor — this is the one piece of raw-number context worth
    # keeping now that the table itself is gone.
    add_text(slide, f"{total_branches} branches  ·  ${total_dep_zones/1e6:.0f}M total deposits",
              6.0, 4.86, 3.8, 0.22, size=8, italic=True, color=GRAY3, align=PP_ALIGN.LEFT)




def build_branches(prs, d, narr, logo_bytes, page_num=2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_num, "PRIORITY MARKETS", logo_bytes)
    add_narrative(slide, narr["priority"], 0.14)

    for i, b in enumerate(d["branchList"][:5]):
        by = 0.14 + i*1.06
        zc  = ZONE_C.get(b["zone"], ANALYZE)
        zbg = ZONE_L.get(b["zone"], ANALYZE_L)

        add_rect(slide, 6.22, by, 3.6, 0.9, GRAY1, GRAY2, Pt(0.4))
        add_rect(slide, 6.22, by, 0.06, 0.9, zc)

        # Rank badge
        add_rect(slide, 6.34, by+0.26, 0.34, 0.34, zc)
        add_text(slide, str(i+1), 6.34, by+0.26, 0.34, 0.34,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, b["name"], 6.76, by+0.06, 2.18, 0.28, size=9.5, bold=True, color=NAVY)
        add_text(slide, b["city"], 6.76, by+0.34, 2.18, 0.18, size=8, color=GRAY3)
        add_text(slide, f"{b['dep']}  ·  {b['yoy']}% YoY", 6.76, by+0.53, 2.18, 0.22, size=9, color=NAVY)

        # Zone pill
        add_rect(slide, 9.0, by+0.30, 0.72, 0.24, zbg, zc, Pt(0.5))
        add_text(slide, b["zone"], 9.0, by+0.30, 0.72, 0.24,
                 size=7, bold=True, color=zc, align=PP_ALIGN.CENTER)


def build_financial(prs, d, narr, logo_bytes, page_num=3):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_num, "FINANCIAL HEALTH", logo_bytes)
    add_narrative(slide, narr["financial"], 0.14)

    # Column headers
    cols = [(6.22, 1.12, "", PP_ALIGN.LEFT),
            (7.38, 1.26, "VALUE",     PP_ALIGN.CENTER),
            (8.68, 1.0,  "BENCHMARK", PP_ALIGN.CENTER)]
    for cx, cw, cl, ca in cols:
        add_rect(slide, cx, 0.14, cw, 0.30, NAVY)
        add_text(slide, cl, cx, 0.14, cw, 0.30, size=7.5, bold=True, color=WHITE, align=ca)

    for i, m in enumerate(d["metrics"]):
        my = 0.48 + i*0.64
        bg = GRAY1 if i%2==0 else WHITE
        add_rect(slide, 6.22, my, 3.56, 0.58, bg, GRAY2, Pt(0.3))
        add_text(slide, m["label"], 6.30, my+0.13, 1.0, 0.28, size=9.5, color=NAVY)
        add_text(slide, m["value"], 7.38, my+0.10, 1.26, 0.32, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, m["bench"], 8.68, my+0.13, 1.0,  0.28, size=9, italic=True, color=GRAY3, align=PP_ALIGN.CENTER)

        sc = TEAL if m["ok"] else AMBER
        add_rect(slide, 9.74, my+0.13, 0.30, 0.30, sc)
        add_text(slide, "✓" if m["ok"] else "!", 9.74, my+0.13, 0.30, 0.30,
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if d.get("competitor"):
        add_rect(slide, 6.22, 4.88, 3.56, 0.34, JUSTIFY_L, JUSTIFY, Pt(0.5))
        add_text(slide,
            f"⚠  Primary Competitor  ·  {d['competitor']['branches']} overlap markets"
            f"  ·  Avg branch vulnerability {d['competitor']['vuln']}/100",
            6.32, 4.90, 3.36, 0.28, size=8.5, bold=True, color=JUSTIFY)

    if d.get("brokered"):
        b = d["brokered"]
        # Brokered callout box — amber alert below metrics
        brok_y = 4.54 if d.get("competitor") else 4.88
        add_rect(slide, 6.22, brok_y, 3.56, 0.62, rgb("FFF8EC"), AMBER, Pt(0.5))
        add_text(slide,
            f"BROKERED DEPOSIT EXPOSURE",
            6.32, brok_y + 0.06, 3.36, 0.20,
            size=7.5, bold=True, color=AMBER)
        add_text(slide,
            f"{b['brokered_pct']}% of deposits ({b['brokered_M']:.0f}M) sourced from brokers"
            f"  ·  Direct savings base: ${b['savings_M']:.0f}M",
            6.32, brok_y + 0.26, 3.36, 0.20,
            size=8, color=NAVY)
        add_text(slide,
            f"Savings→CD funnel strategy could reduce funding cost by converting"
            f" brokered volume to direct customer deposits.",
            6.32, brok_y + 0.42, 3.36, 0.18,
            size=7.5, italic=True, color=GRAY3)


def build_gap(prs, d, narr, page_num=4):
    """Dark slide — no chrome, uses navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = NAVY

    # Left teal stripe only
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)

    # Giant gap number
    add_text(slide, d["gap"], 0.28, 0.16, 5.2, 1.86,
             size=96, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    add_text(slide, "GAP VS PEER AVERAGE", 0.28, 2.08, 5.2, 0.34,
             size=13, bold=True, color=WHITE)
    add_text(slide, d["gapSubtitle"], 0.28, 2.50, 5.2, 0.26,
             size=9.5, italic=True, color=rgb("4A6A8A"))

    # 3 stat tiles
    tile_c = rgb("F87171") if d["gapNeg"] else TEAL
    tiles = [
        (f"{d['bankYoY']}%", f"{d['bankName'].upper()} YoY", tile_c),
        (f"{d['peerYoY']}%", "PEER AVG",      GRAY3),
        (d["gap"],           "GAP",            AMBER),
    ]
    for i, (val, lbl, c) in enumerate(tiles):
        tx = 0.28 + i*1.78
        add_rect(slide, tx, 2.96, 1.62, 1.12, rgb("162436"), rgb("1E3A5F"), Pt(0.5))
        add_rect(slide, tx, 2.96, 1.62, 0.06, c)
        # First tile carries the full bank name — give it a smaller font
        # and 2-line-capable box instead of truncating mid-name.
        lbl_size = 6 if i == 0 and len(d["bankName"]) > 14 else 7
        add_text(slide, val, tx, 3.02, 1.62, 0.52, size=20, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, tx+0.04, 3.58, 1.54, 0.42, size=lbl_size, bold=True,
                 color=rgb("3A5A7A"), align=PP_ALIGN.CENTER)

    # ── NATIVE VECTOR BAR CHART ──────────────────────────────────
    # Full bank name — category axis wraps naturally within column width,
    # no truncation needed (truncating to ~20 chars was cutting names like
    # "Citizens Independent Bank" mid-name even at a word boundary).
    chart_data = ChartData()
    chart_data.categories = [d["bankName"], "Peer Avg"]
    chart_data.add_series("Deposit YoY %", (float(d["bankYoY"]), float(d["peerYoY"])))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.4), Inches(0.15), Inches(4.4), Inches(5.0),
        chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = False
    chart.has_title  = False

    # Bar colors
    bar_colors = [rgb("A32D2D") if d["gapNeg"] else TEAL, GRAY3]
    for i, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = bar_colors[i]
        point.format.line.fill.background()

    # Style axes
    va = chart.value_axis
    va.tick_labels.font.color.rgb = GRAY3
    va.tick_labels.font.size = Pt(11)

    ca = chart.category_axis
    ca.tick_labels.font.color.rgb = WHITE
    ca.tick_labels.font.size = Pt(13) if len(d["bankName"]) <= 16 else Pt(10)
    ca.tick_labels.font.bold = True

    # Footer
    add_text(slide,
        f"Verlocity Princeton Partners Group   ·   BMAP Intelligence   ·   {d['bankName']}",
        0.28, 5.30, 9.5, 0.22, size=7.5, color=rgb("8DA3BC"))
    add_text(slide, str(page_num), 9.50, 5.30, 0.38, 0.22, size=9, color=rgb("8DA3BC"), align=PP_ALIGN.RIGHT)


def build_next_steps(prs, d, narr, logo_bytes, page_num=6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_chrome(slide, page_num, "THE VERLOCITY PLATFORM", logo_bytes)
    add_narrative(slide, narr["nextsteps"], 0.14)

    # Product labels above cards — all four presented as live capabilities
    phase_labels = ["BMAP", "AUDIENCE/MEDIA", "OMNIBRANCH", "PREDICTIVE ROI"]
    phase_colors = [TEAL, ANALYZE, AMBER, NAVY]
    ac_colors = [TEAL, ANALYZE, AMBER, NAVY]

    for i, action in enumerate(d["actions"]):
        ay = 0.14 + i*1.30
        ac = ac_colors[i]

        add_rect(slide, 6.22, ay, 3.6, 1.14, GRAY1, GRAY2, Pt(0.4))
        add_rect(slide, 6.22, ay, 0.06, 1.14, ac)
        add_rect(slide, 6.34, ay+0.14, 0.34, 0.34, ac)
        add_text(slide, str(i+1).zfill(2), 6.34, ay+0.14, 0.34, 0.34,
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Phase pill
        add_rect(slide, 6.76, ay+0.10, 0.90, 0.16, phase_colors[i])
        add_text(slide, phase_labels[i], 6.76, ay+0.10, 0.90, 0.16,
                 size=5.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, action["title"], 6.76, ay+0.28, 2.98, 0.26,
                 size=10, bold=True, color=NAVY)
        add_text(slide, action["body"],  6.76, ay+0.56, 2.98, 0.52,
                 size=8, color=GRAY3)


# ═══════════════════════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════════════════════

def _build_branch_list(br, sf):
    """
    Urgency mix: top Invest branches (upside) + critical Defend/Justify
    branches by deposit size (risk). Scores redacted — zone pill only.
    Max 5 cards. Names up to 28 chars.
    """
    invest_br  = sorted(
        [b for b in br if b.get("opportunity_zone") == "Invest" and sf(b.get("latest_dep")) >= 5e6],
        key=lambda b: sf(b.get("opportunity_score")), reverse=True
    )[:3]
    risk_br = sorted(
        [b for b in br if b.get("opportunity_zone") in ("Defend", "Justify")],
        key=lambda b: sf(b.get("latest_dep")), reverse=True
    )[:2]

    # Merge, dedupe, cap at 5
    seen = set()
    merged = []
    for b in invest_br + risk_br:
        uid = b.get("uninumbr") or b.get("namebr")
        if uid not in seen:
            seen.add(uid)
            merged.append(b)
    merged = merged[:5]

    return [
        {
            "name":  b["namebr"].split("--")[-1].strip()[:28],
            "city":  f"{b.get('citybr','')}, {b.get('stalpbr','')}",
            "dep":   f"${sf(b.get('latest_dep'))/1e6:.0f}M",
            "yoy":   f"{sf(b.get('yoy_deposits'))*100:+.1f}",
            "zone":  b.get("opportunity_zone", ""),
        }
        for b in merged
    ]



# ═══════════════════════════════════════════════════════════════
# PERSONA SLIDE
# ═══════════════════════════════════════════════════════════════

def build_persona_slide(prs, personas, bank_name, logo_bytes, page_num=5):
    """
    Slide: Top 3 Audience Personas — bridge to AudienceFinder.
    Layout: dark navy header, 3 side-by-side persona cards, audiencefinder CTA bar.
    """
    GREEN   = rgb("2ECC71")
    TEAL    = rgb("028090")
    CARD_BG = rgb("F7F8FA")
    CARD_BD = rgb("DDE3EA")
    ACCENT  = rgb("00A896")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left nav bar (consistent with other slides)
    add_rect(slide, 0, 0, 0.28, 5.625, NAVY)

    # Header background
    add_rect(slide, 0.28, 0, 9.72, 1.05, NAVY)

    # Slide title
    add_text(slide, "YOUR MARKET IS TELLING YOU WHO TO TALK TO",
             0.42, 0.08, 7.5, 0.38, size=13, bold=True, color=WHITE)
    add_text(slide, f"Top 3 audience segments across Invest + Analyze branches · {bank_name}",
             0.42, 0.52, 7.5, 0.30, size=8.5, color=rgb("8BAECC"), italic=True)

    # AudienceFinder label top right (plain text, no button)
    add_text(slide, "POWERED BY AUDIENCEFINDER", 7.50, 0.18, 2.30, 0.28,
             size=7, bold=True, color=rgb("00A896"), align=PP_ALIGN.RIGHT)

    # 3 persona cards side by side
    card_w = 2.90
    card_h = 3.60
    card_y = 1.15
    gaps   = [0.42, 3.46, 6.50]  # x positions for 3 cards

    persona_colors = [rgb("1A6B8A"), rgb("1A7A5E"), rgb("6B3FA0")]  # teal, green, purple
    icons = ["◎", "◉", "●"]  # simple circle icons

    for i, (cx, p) in enumerate(zip(gaps, personas[:3])):
        accent_c = persona_colors[i]

        # Card background
        add_rect(slide, cx, card_y, card_w, card_h, CARD_BG, CARD_BD, Pt(0.75))

        # Card top accent strip
        add_rect(slide, cx, card_y, card_w, 0.08, accent_c)

        # Persona number badge
        add_rect(slide, cx + 0.12, card_y + 0.14, 0.30, 0.30, accent_c)
        add_text(slide, str(i+1), cx + 0.12, card_y + 0.14, 0.30, 0.30,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Persona name
        add_text(slide, p.get("name","—"), cx + 0.50, card_y + 0.14, card_w - 0.62, 0.32,
                 size=10.5, bold=True, color=NAVY)

        # Demographic pills row
        age_str    = p.get("age","")
        income_str = p.get("income","")
        occ_str    = p.get("occupation","")
        demo_line  = f"{age_str}  ·  {income_str}"
        add_text(slide, demo_line, cx + 0.12, card_y + 0.52, card_w - 0.24, 0.22,
                 size=8, color=rgb("446688"), bold=False)

        # Occupation
        add_text(slide, occ_str[:45], cx + 0.12, card_y + 0.74, card_w - 0.24, 0.20,
                 size=7.5, color=GRAY3, italic=True)

        # Divider
        add_rect(slide, cx + 0.12, card_y + 1.00, card_w - 0.24, 0.02, CARD_BD)

        # Insight
        add_text(slide, "INSIGHT", cx + 0.12, card_y + 1.10, card_w - 0.24, 0.18,
                 size=6.5, bold=True, color=accent_c)
        add_text(slide, p.get("insight",""), cx + 0.12, card_y + 1.28, card_w - 0.24, 0.68,
                 size=8, color=NAVY)

        # Banking moment
        add_text(slide, "BANKING MOMENT", cx + 0.12, card_y + 2.02, card_w - 0.24, 0.18,
                 size=6.5, bold=True, color=accent_c)
        add_text(slide, p.get("moment",""), cx + 0.12, card_y + 2.20, card_w - 0.24, 0.42,
                 size=8, color=NAVY)

        # Why now
        add_rect(slide, cx + 0.12, card_y + 2.68, card_w - 0.24, 0.02, CARD_BD)
        add_text(slide, "WHY NOW", cx + 0.12, card_y + 2.76, card_w - 0.24, 0.16,
                 size=6.5, bold=True, color=GRAY3)
        add_text(slide, p.get("why_now",""), cx + 0.12, card_y + 2.92, card_w - 0.24, 0.50,
                 size=7.5, color=GRAY3, italic=True)

    # Bottom CTA bar
    add_rect(slide, 0.28, 4.88, 9.52, 0.56, NAVY)
    add_text(slide,
             "AudienceFinder builds precision digital audiences around these 3 segments — "
             "reach them before your competitors do.",
             0.42, 4.88, 9.20, 0.56, size=8.5, color=WHITE, italic=True, valign="center")

    # Logo bottom left
    if logo_bytes:
        try:
            import io
            logo_stream = io.BytesIO(logo_bytes)
            slide.shapes.add_picture(logo_stream, Inches(0.35), Inches(4.90),
                                     width=Inches(0.90), height=Inches(0.45))
        except Exception:
            pass

    # Page number
    add_text(slide, str(page_num), 9.50, 5.28, 0.38, 0.20, size=9, color=GRAY3, align=PP_ALIGN.RIGHT)

    return slide

def build_deck(data, logo_bytes):
    rows = data["rows"]; fin = data["fin"]; tgt = data["tgt"]
    br   = data["br"];   bankName = data["bankName"]

    def sf(v, default=0):
        """Safe float — handles None, empty string, missing."""
        try:
            return float(v) if v is not None else default
        except (TypeError, ValueError):
            return default

    tot      = sum(sf(r.get("latest_dep")) for r in rows)
    avg      = lambda v: sum(sf(r.get(v)) for r in rows) / max(len(rows),1)
    invest   = sum(1 for r in rows if r.get("opportunity_zone")=="Invest")
    analyze  = sum(1 for r in rows if r.get("opportunity_zone")=="Analyze")
    defend   = sum(1 for r in rows if r.get("opportunity_zone")=="Defend")
    justify  = sum(1 for r in rows if r.get("opportunity_zone")=="Justify")
    dep_invest  = sum(sf(r.get("latest_dep")) for r in rows if r.get("opportunity_zone")=="Invest")
    dep_analyze = sum(sf(r.get("latest_dep")) for r in rows if r.get("opportunity_zone")=="Analyze")
    dep_defend  = sum(sf(r.get("latest_dep")) for r in rows if r.get("opportunity_zone")=="Defend")
    dep_justify = sum(sf(r.get("latest_dep")) for r in rows if r.get("opportunity_zone")=="Justify")
    bankYoY  = avg("yoy_deposits")*100
    compYoY  = avg("avg_comp_yoy")*100
    gap      = bankYoY - compYoY
    avgScore = avg("opportunity_score")

    top_br   = sorted(br, key=lambda b: sf(b.get("opportunity_score")), reverse=True)
    just_top = sorted([b for b in br if b.get("opportunity_zone")=="Justify"],
                      key=lambda b: sf(b.get("latest_dep")), reverse=True)
    tier1    = [b for b in br if b.get("campaign") in ["Aggressive Acquisition","Urgent Competitive Push","Capitalize","Turnaround","Grow Share","Competitive Defense"]][:2]

    narr = get_narratives(data)

    D = {
        "bankName":    bankName,
        "date":        datetime.now().strftime("%B %Y"),
        "branchCount": str(len(rows)),
        "deposits":    f"${tot/1e9:.1f}B",
        "avgScore":    f"{avgScore:.0f}",
        "depositYoY":  f"{bankYoY:+.1f}%",
        "gap":         f"{gap:+.1f}pp",
        "bankYoY":     f"{bankYoY:.1f}",
        "peerYoY":     f"{compYoY:.1f}",
        "gapNeg":      gap < 0,
        "gapSubtitle": f"Deposit growth vs. peer average — {fin.get('period','Q4 2025')}",
        "invest":  invest,  "analyze": analyze,
        "defend":  defend,  "justify": justify,
        "depInvest":  dep_invest,  "depAnalyze": dep_analyze,
        "depDefend":  dep_defend,  "depJustify": dep_justify,
        "branchList": _build_branch_list(br, sf),
        "metrics": [
            {"label":"ROA",           "value":f"{sf(fin.get('roa')):.2f}%",              "bench":">1.0%",    "ok": sf(fin.get("roa"))>=1},
            {"label":"NIM",           "value":f"{sf(fin.get('nim')):.2f}%",              "bench":"2.5–3.5%", "ok": 2.5<=sf(fin.get("nim"))<=4},
            {"label":"Efficiency",    "value":f"{sf(fin.get('efficiency_ratio')):.2f}%", "bench":"<60%",     "ok": 0<sf(fin.get("efficiency_ratio"))<60},
            {"label":"Net Income YoY","value":f"{sf(fin.get('net_income_yoy_pct')):+.1f}%",               "bench":">0%",      "ok": sf(fin.get("net_income_yoy_pct"))>0},
            {"label":"Deposit YoY",   "value":f"{bankYoY:+.1f}%",                        "bench":">2%",      "ok": bankYoY>=2},
            {"label":"Cost of Funds", "value":f"{sf(fin.get('cost_of_funds_pct')):.2f}%","bench":"<2%",      "ok": 0<sf(fin.get("cost_of_funds_pct"))<2},
            {"label":"Tier 1 Capital","value":f"{sf(fin.get('tier1_capital_pct')):.2f}%","bench":">8%",      "ok": sf(fin.get("tier1_capital_pct"))>=8},
        ],
        "competitor": {
            "branches": tgt["branches_in_radius"],
            "yoy":      f"{sf(tgt.get('avg_yoy_pct')):.1f}",
            "vuln":     f"{sf(tgt.get('avg_vuln_score')):.0f}",
        } if tgt else None,
        "brokered": data.get("brokered"),
        "actions": [
            {
                "title": "BMAP — Market Truth",
                "body":  f"Full branch scoring across all {len(rows)} locations. Decision-quality clarity on where to invest, defend, and exit. The data foundation every other step builds on.",
            },
            {
                "title": "AudienceFinder & MediaPredict",
                "body":  f"Rate-sensitive depositor campaigns targeted to your top {invest} Invest-zone markets. Precision media that general spend can't replicate — savings openers converted into sticky CD relationships.",
            },
            {
                "title": "Omnibranch",
                "body":  "Unified branch-level orchestration — connects deposit data, market signals, and campaign execution into one operating view across the network.",
            },
            {
                "title": "Predictive ROI",
                "body":  "Forecasts the return on each marketing dollar before it's spent, then tracks performance against that forecast — accountability built into the system, not added after.",
            },
        ],
    }

    print(f"  Building slides...")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    # Ensure blank layout exists
    while len(prs.slide_layouts) < 7:
        prs.slide_layouts.add_slide()

    build_cover(prs, D, logo_bytes)
    build_network(prs, D, narr, logo_bytes, page_num=1)
    build_branches(prs, D, narr, logo_bytes, page_num=2)
    build_financial(prs, D, narr, logo_bytes, page_num=3)
    build_gap(prs, D, narr, page_num=4)
    # Persona slide — before next steps
    personas = data.get("personas")
    next_page = 5
    if personas and len(personas) >= 1:
        print(f"  Adding persona slide ({len(personas)} personas)...")
        build_persona_slide(prs, personas, bankName, logo_bytes, page_num=5)
        next_page = 6
    else:
        print("  Skipping persona slide — no personas available")
    build_next_steps(prs, D, narr, logo_bytes, page_num=next_page)

    return prs


def save_deck(prs, bank_name, out_dir=OUT_DIR):
    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in bank_name).strip()
    date = datetime.now().strftime("%Y%m%d")
    fname = out_dir / f"BMAP_Snapshot_{safe}_{date}.pptx"
    prs.save(str(fname))
    return fname


# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════

def run_single(ik, name_hint=None):
    print(f"\n{'='*55}")
    print(f"  BMAP Snapshot — {name_hint or ik}")
    print(f"{'='*55}")
    data = fetch_bank_data(ik)
    if name_hint:
        data["bankName"] = name_hint
    # Fetch or generate personas
    personas = fetch_or_generate_personas(
        ik, data["bankName"], data.get("br", []), data)
    data["personas"] = personas
    logo = fetch_logo()
    prs  = build_deck(data, logo)
    path = save_deck(prs, data["bankName"])
    print(f"\n  ✓  Saved: {path}\n")
    return path


def run_batch(csv_path):
    banks = []
    with open(csv_path, newline="") as f:
        for row in csv.reader(f):
            if not row or row[0].startswith("#") or row[0].lower()=="inst_key":
                continue
            ik   = row[0].strip()
            name = row[1].strip() if len(row)>1 else ""
            banks.append((ik, name))

    print(f"\n  Batch: {len(banks)} banks from {csv_path}")
    logo = fetch_logo()
    out_dir = OUT_DIR / f"BMAP_Batch_{datetime.now().strftime('%Y%m%d_%H%M')}"
    out_dir.mkdir(parents=True, exist_ok=True)

    results = []
    for i, (ik, name) in enumerate(banks):
        print(f"\n[{i+1}/{len(banks)}] {name or ik}")
        try:
            data = fetch_bank_data(ik)
            if name: data["bankName"] = name
            prs  = build_deck(data, logo)
            path = save_deck(prs, data["bankName"], out_dir)
            print(f"  ✓  {path.name}")
            results.append({"bank": data["bankName"], "status": "ok", "file": str(path)})
        except Exception as e:
            print(f"  ✗  FAILED: {e}")
            results.append({"bank": name or ik, "status": "error", "error": str(e)})
        if i < len(banks)-1:
            time.sleep(0.5)

    ok  = sum(1 for r in results if r["status"]=="ok")
    err = sum(1 for r in results if r["status"]=="error")
    print(f"\n{'='*55}")
    print(f"  Batch complete — {ok} decks saved, {err} failed")
    print(f"  Output folder: {out_dir}")
    print(f"{'='*55}\n")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="BMAP Snapshot Deck Builder — generates branded PPTX from Supabase data"
    )
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--inst_key", help="Single bank inst_key (e.g. bank_123)")
    group.add_argument("--csv",      help="CSV file with inst_key[,name] per line")
    parser.add_argument("--name",    help="Override bank name (single mode only)")
    parser.add_argument("--out",     help="Output directory (default: current folder)", default=".")
    args = parser.parse_args()

    OUT_DIR = Path(args.out)
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    if not ANTH_KEY:
        print("\n  ⚠  Set ANTHROPIC_API_KEY env var for AI narratives.")
        print("     export ANTHROPIC_API_KEY=sk-ant-...")
        print("     Continuing with empty narratives...\n")

    if args.inst_key:
        run_single(args.inst_key, args.name)
    else:
        run_batch(args.csv)
